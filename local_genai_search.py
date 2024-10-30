import os
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import PyPDF2
import docx
from pptx import Presentation
import json
import streamlit as st
import re
import ollama
from streamlit_lottie import st_lottie
import requests

print("Starting the application...")

# Global variables
model = SentenceTransformer('sentence-transformers/paraphrase-multilingual-mpnet-base-v2')
dimension = 768
index = faiss.IndexFlatIP(dimension)
metadata = []
documents_path = ""  # Add this line to define documents_path globally

print(f"Initialized model and FAISS index with dimension {dimension}")

# Document reading functions
def read_pdf(file_path):
    print(f"Reading PDF: {file_path}")
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        return ' '.join([page.extract_text() for page in reader.pages])

def read_docx(file_path):
    print(f"Reading DOCX: {file_path}")
    doc = docx.Document(file_path)
    return ' '.join([para.text for para in doc.paragraphs])

def read_pptx(file_path):
    print(f"Reading PPTX: {file_path}")
    prs = Presentation(file_path)
    return ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])

def chunk_text(text, chunk_size=500, overlap=50):
    print(f"Chunking text of length {len(text)} with chunk size {chunk_size} and overlap {overlap}")
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = ' '.join(words[i:i + chunk_size])
        chunks.append(chunk)
    print(f"Created {len(chunks)} chunks")
    return chunks

# Indexing function
def index_documents(directory):
    print(f"Indexing documents in directory: {directory}")
    global metadata, index
    metadata = []  # Reset metadata
    documents = []
    index = faiss.IndexFlatIP(dimension)  # Reset the index
    
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            print(f"Processing file: {file_path}")
            content = ""
            
            if file.lower().endswith('.pdf'):
                content = read_pdf(file_path)
            elif file.lower().endswith('.docx'):
                content = read_docx(file_path)
            elif file.lower().endswith('.pptx'):
                content = read_pptx(file_path)
            elif file.lower().endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            
            if content:
                chunks = chunk_text(content)
                for i, chunk in enumerate(chunks):
                    documents.append(chunk)
                    metadata.append({"path": file_path, "chunk_id": i})
    
    print(f"Encoding {len(documents)} document chunks")
    embeddings = model.encode(documents)
    print(f"Adding embeddings to FAISS index")
    index.add(np.array(embeddings))
    
    # Save index and metadata
    print("Saving FAISS index and metadata")
    faiss.write_index(index, "document_index.faiss")
    with open("metadata.json", "w") as f:
        json.dump(metadata, f)
    
    print(f"Indexed {len(documents)} document chunks.")

# Function to read document chunk
def read_document_chunk(file_path, chunk_id):
    global documents_path 
    print(f"Reading document chunk: {file_path}, chunk_id: {chunk_id}")
    content = ""
    
    
    if not os.path.exists(file_path):
        base_name = os.path.basename(file_path)
        new_path = os.path.join(documents_path, base_name)
        if os.path.exists(new_path):
            file_path = new_path
        else:
            print(f"File not found: {file_path}")
            return f"[Content not available for {base_name}]"
    
    if file_path.endswith('.pdf'):
        content = read_pdf(file_path)
    elif file_path.endswith('.docx'):
        content = read_docx(file_path)
    elif file_path.endswith('.pptx'):
        content = read_pptx(file_path)
    elif file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    
    chunks = chunk_text(content)
    return chunks[chunk_id] if chunk_id < len(chunks) else ""

# Search function
def semantic_search(query, k=10):
    print(f"Performing semantic search for query: '{query}', k={k}")
    query_vector = model.encode([query])[0]
    distances, indices = index.search(np.array([query_vector]), k)
    
    results = []
    for i, idx in enumerate(indices[0]):
        meta = metadata[idx]
        content = read_document_chunk(meta["path"], meta["chunk_id"])
        results.append({
            "id": int(idx),
            "path": meta["path"],
            "content": content,
            "score": float(distances[0][i])
        })
    
    print(f"Found {len(results)} search results")
    return results

# Answer generation function
def generate_answer(query, context):
    print(f"Generating answer for query: '{query}'")
    prompt = f"""Systemanweisung: Du bist ein hilfsbereiter Assistent. Antworte immer auf Deutsch, 
unabhängig von der Sprache der Frage. Beantworten Sie die Frage des Benutzers anhand der im Kontext angegebenen Dokumente. Im Kontext stehen Dokumente, die eine Antwort enthalten sollen. Bitte geben Sie immer die Dokument-ID (in eckigen Klammern, z.B. [0],[1]) des Dokuments an, das für eine Behauptung verwendet wurde. Verwenden Sie so viele Zitate und Dokumente wie nötig, um die Frage zu beantworten.

Context: {context} 

Question: {query}

Answer:"""

    print("Sending prompt to Ollama")
    response = ollama.generate(model='llama3.2', prompt=prompt)
    print("Received response from Ollama")
    return response['response']

def load_lottieurl(url: str):
    r = requests.get(url)
    if r.status_code != 200:
        return None
    return r.json()

# Streamlit UI
def main():
    global documents_path, index, metadata
    print("Starting Streamlit UI")
    
    # Page config
    st.set_page_config(page_title="Local GenAI Search", page_icon="🔍", layout="wide")
    
    # Custom CSS
    st.markdown("""
    <style>
    .big-font {
        font-size:30px !important;
        font-weight: bold;
        color: #1E90FF;
    }
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        font-size: 18px;
        padding: 10px 24px;
        border-radius: 12px;
        border: none;
    }
    .stTextInput>div>div>input {
        font-size: 16px;
    }
    </style>
    """, unsafe_allow_html=True)

    # Title and animation
    col1, col2 = st.columns([2, 1])
    with col1:
        st.markdown('<p class="big-font">Local GenAI Search 🔍</p>', unsafe_allow_html=True)
        st.write("Explore your documents with the power of AI!")
    with col2:
        lottie_url = "https://assets5.lottiefiles.com/packages/lf20_fcfjwiyb.json"
        lottie_json = load_lottieurl(lottie_url)
        st_lottie(lottie_json, height=150, key="coding")

    # Input for documents path
    new_documents_path = st.text_input("📁 Enter the path to your documents folder:", "")
    if new_documents_path and new_documents_path != documents_path:
        documents_path = os.path.abspath(new_documents_path)
        # Reset index and metadata
        index = faiss.IndexFlatIP(dimension)
        metadata = []
        if st.button("🚀 Index Documents"):
            with st.spinner("Indexing documents... This may take a while."):
                print(f"Indexing documents in {documents_path}")
                index_documents(documents_path)
            st.success("✅ Indexing complete!")
            st.rerun()  # Use st.rerun() instead of st.experimental_rerun()
    
    # Load index and metadata if they exist
    if os.path.exists("document_index.faiss") and os.path.exists("metadata.json"):
        if len(metadata) == 0:
            print("Loading FAISS index and metadata")
            index = faiss.read_index("document_index.faiss")
            with open("metadata.json", "r") as f:
                metadata = json.load(f)
            print(f"Loaded index with {index.ntotal} vectors and {len(metadata)} metadata entries")
    else:
        st.warning("⚠️ Documents are not indexed. Please enter a folder path and click 'Index Documents'.")
    
    st.markdown("---")
    st.markdown("## Ask a Question")
    question = st.text_input("🤔 What would you like to know about your documents?", "")

    if st.button("🔍 Search and Answer"):
        if question:
            with st.spinner("Searching and generating answer..."):
                print(f"User asked: '{question}'")
                
                # Perform semantic search
                search_results = semantic_search(question)
                
                # Prepare context for answer generation
                context = "\n\n".join([f"{i}: {result['content']}" for i, result in enumerate(search_results)])
                # Generate answer
                answer = generate_answer(question, context)
                
                st.markdown("### 🤖 AI Answer:")
                st.markdown(answer)
                
                # Display referenced documents
                st.markdown("### 📚 Referenced Documents:")
                rege = re.compile(r"\[Document\s+[0-9]+\]|\[[0-9]+\]")
                referenced_ids = [int(s) for s in re.findall(r'\b\d+\b', ' '.join(rege.findall(answer)))]
                
                print(f"Displaying {len(referenced_ids)} referenced documents")
                for doc_id in referenced_ids:
                    doc = search_results[doc_id]
                    full_path = os.path.join(documents_path, os.path.basename(doc['path']))
                    if os.path.exists(full_path):
                        with st.expander(f"📄 Document {doc_id} - {os.path.basename(doc['path'])}"):
                            st.write(doc['content'])
                            with open(full_path, 'rb') as f:
                                st.download_button("⬇️ Download file", f, file_name=os.path.basename(doc['path']))
                    else:
                        st.warning(f"⚠️ File not found: {full_path}")
        else:
            st.warning("⚠️ Please enter a question before clicking 'Search and Answer'.")

if __name__ == "__main__":
    main()
    print("Application finished")
